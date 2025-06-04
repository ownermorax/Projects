import os
import io
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from PIL import Image
import pdf2image
import pytesseract
from pdf2docx import Converter
from pptx import Presentation
import img2pdf
from fpdf import FPDF
from docx import Document
import comtypes.client
import pymupdf
import fitz

class FileConverterApp:
    def __init__(self, root):
        self.root = root
        self.root.title("Универсальный конвертер файлов")
        self.root.geometry("900x700")
        self.root.resizable(False, False)
        
        self.bg_color = "#000000" 
        self.fg_color = "#00FF00" 
        self.button_color = "#003300" 
        self.text_color = "#00FF00"
        self.highlight_color = "#006600" 
        
        self.style = ttk.Style()
        self.style.configure("TFrame", background=self.bg_color)
        self.style.configure("TLabel", background=self.bg_color, foreground=self.fg_color, font=('Arial', 10))
        self.style.configure("TButton", background=self.button_color, foreground=self.fg_color, font=('Arial', 10))
        self.style.configure("TCombobox", fieldbackground=self.bg_color, foreground=self.fg_color, background=self.bg_color)
        
        self.supported_formats = {
            "PDF": ["TXT", "DOCX", "PNG", "JPG"],
            "TXT": ["PDF", "DOCX"],
            "DOCX": ["PDF", "TXT"],
            "PPTX": ["PDF"],
            "PNG": ["JPG", "GIF", "PDF", "ICO"],
            "JPG": ["PNG", "GIF", "PDF", "ICO"],
            "GIF": ["PNG", "JPG", "PDF"],
            "ICO": ["PNG", "JPG"]
        }
        
        self.conversion_list = """
Поддерживаемые конвертации:
-----------------------------
PDF → TXT, DOCX, PNG, JPG
TXT → PDF, DOCX
DOCX → PDF, TXT
PPTX → PDF
PNG ↔ JPG ↔ GIF ↔ PDF ↔ ICO
JPG ↔ PNG ↔ GIF ↔ PDF ↔ ICO
ICO ↔ PNG ↔ JPG
-----------------------------
"""
        
        self.create_widgets()
    
    def create_widgets(self):
        self.main_frame = ttk.Frame(self.root)
        self.main_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=20)
        
        self.title_label = ttk.Label(
            self.main_frame, 
            text="Универсальный конвертер файлов", 
            font=('Arial', 16, 'bold')
        )
        self.title_label.pack(pady=(0, 10))
        
        self.info_frame = ttk.Frame(self.main_frame, style="TFrame")
        self.info_frame.pack(fill=tk.X, pady=5)
        
        self.info_label = tk.Label(
            self.info_frame,
            text=self.conversion_list,
            bg=self.bg_color,
            fg=self.fg_color,
            font=('Courier New', 9),
            justify=tk.LEFT
        )
        self.info_label.pack(anchor=tk.W)
        
        self.file_frame = ttk.Frame(self.main_frame)
        self.file_frame.pack(fill=tk.X, pady=10)
        
        self.file_label = ttk.Label(self.file_frame, text="Выберите файл:")
        self.file_label.pack(side=tk.LEFT, padx=(0, 10))
        
        self.file_entry = tk.Entry(
            self.file_frame,
            width=50,
            bg=self.bg_color,
            fg=self.fg_color,
            insertbackground=self.fg_color,
            font=('Arial', 10)
        )
        self.file_entry.pack(side=tk.LEFT, expand=True, fill=tk.X)
        
        self.browse_button = ttk.Button(
            self.file_frame, 
            text="Обзор...", 
            command=self.browse_file
        )
        self.browse_button.pack(side=tk.LEFT, padx=(10, 0))
        
        self.format_frame = ttk.Frame(self.main_frame)
        self.format_frame.pack(fill=tk.X, pady=10)
        
        self.from_label = ttk.Label(self.format_frame, text="Из формата:")
        self.from_label.pack(side=tk.LEFT, padx=(0, 10))
        
        self.from_combobox = ttk.Combobox(
            self.format_frame,
            width=15,
            state="readonly",
            style="TCombobox"
        )
        self.from_combobox.pack(side=tk.LEFT)
        self.from_combobox.bind("<<ComboboxSelected>>", self.update_to_formats)
        
        self.to_label = ttk.Label(self.format_frame, text="В формат:")
        self.to_label.pack(side=tk.LEFT, padx=(10, 10))
        
        self.to_combobox = ttk.Combobox(
            self.format_frame,
            width=15,
            state="readonly",
            style="TCombobox"
        )
        self.to_combobox.pack(side=tk.LEFT)
        
        self.convert_button = ttk.Button(
            self.main_frame, 
            text="Конвертировать", 
            command=self.convert_file,
            style="TButton"
        )
        self.convert_button.pack(pady=15)
        
        self.log_frame = ttk.Frame(self.main_frame)
        self.log_frame.pack(fill=tk.BOTH, expand=True)
        
        self.log_label = ttk.Label(self.log_frame, text="Журнал операций:")
        self.log_label.pack(anchor=tk.W)
        
        self.log_text = tk.Text(
            self.log_frame, 
            height=10, 
            wrap=tk.WORD, 
            bg=self.bg_color, 
            fg=self.fg_color,
            insertbackground=self.fg_color,
            font=('Courier New', 9)
        )
        self.log_text.pack(fill=tk.BOTH, expand=True)
        
        self.scrollbar = ttk.Scrollbar(self.log_text)
        self.scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.log_text.config(yscrollcommand=self.scrollbar.set)
        self.scrollbar.config(command=self.log_text.yview)

    def browse_file(self):
        file_path = filedialog.askopenfilename()
        if file_path:
            self.file_entry.delete(0, tk.END)
            self.file_entry.insert(0, file_path)
            
            file_ext = file_path.split('.')[-1].upper()
            if file_ext in self.supported_formats:
                self.from_combobox.set(file_ext)
                self.update_to_formats()
            else:
                messagebox.showwarning("Предупреждение", "Данный формат файла не поддерживается")
    
    def update_to_formats(self, event=None):
        from_format = self.from_combobox.get()
        if from_format in self.supported_formats:
            self.to_combobox['values'] = self.supported_formats[from_format]
            if self.supported_formats[from_format]:
                self.to_combobox.current(0)
    
    def log_message(self, message):
        self.log_text.insert(tk.END, message + "\n")
        self.log_text.see(tk.END)
        self.root.update()
    
    def convert_file(self):
        input_file = self.file_entry.get()
        if not input_file:
            messagebox.showerror("Ошибка", "Пожалуйста, выберите файл для конвертации")
            return
            
        from_format = self.from_combobox.get()
        to_format = self.to_combobox.get()
        
        if not from_format or not to_format:
            messagebox.showerror("Ошибка", "Пожалуйста, выберите исходный и конечный форматы")
            return
            
        try:
            self.log_message(f"Начало конвертации из {from_format} в {to_format}...")
            output_file = os.path.splitext(input_file)[0] + f"_converted.{to_format.lower()}"
            if from_format == "PDF" and to_format == "TXT":
                self.pdf_to_txt(input_file, output_file)
            elif from_format == "PDF" and to_format == "DOCX":
                self.pdf_to_docx(input_file, output_file)
            elif from_format == "PDF" and to_format in ["PNG", "JPG"]:
                self.pdf_to_image(input_file, output_file, to_format.lower())
            elif from_format == "TXT" and to_format == "PDF":
                self.txt_to_pdf(input_file, output_file)
            elif from_format == "TXT" and to_format == "DOCX":
                self.txt_to_docx(input_file, output_file)
            elif from_format == "DOCX" and to_format == "PDF":
                self.docx_to_pdf(input_file, output_file)
            elif from_format == "DOCX" and to_format == "TXT":
                self.docx_to_txt(input_file, output_file)
            elif from_format == "PPTX" and to_format == "PDF":
                self.pptx_to_pdf(input_file, output_file)
            elif from_format in ["PNG", "JPG", "GIF"] and to_format in ["PNG", "JPG", "GIF"]:
                self.convert_image(input_file, output_file, to_format.lower())
            elif from_format in ["PNG", "JPG", "GIF"] and to_format == "PDF":
                self.image_to_pdf(input_file, output_file)
            elif from_format in ["PNG", "JPG"] and to_format == "ICO":
                self.image_to_ico(input_file, output_file)
            elif from_format == "ICO" and to_format in ["PNG", "JPG"]:
                self.ico_to_image(input_file, output_file, to_format.lower())
            else:
                messagebox.showerror("Ошибка", "Данное преобразование не поддерживается")
                return
            self.log_message(f"Конвертация завершена. Результат сохранен в: {output_file}")
            messagebox.showinfo("Успех", "Файл успешно сконвертирован!")
            
        except Exception as e:
            self.log_message(f"Ошибка при конвертации: {str(e)}")
            messagebox.showerror("Ошибка", f"Произошла ошибка при конвертации: {str(e)}")
    
    def pdf_to_txt(self, input_file, output_file):
        self.log_message("Извлечение текста из PDF...")
        try:
            from pdfminer.high_level import extract_text
            text = extract_text(input_file)
            with open(output_file, "w", encoding="utf-8") as f:
                f.write(text)
        except Exception as e:
            self.log_message(f"Ошибка при обработке PDF: {str(e)}")
            raise
    
    def pdf_to_docx(self, input_file, output_file):
        self.log_message("Конвертация PDF в DOCX...")
        try:
            cv = Converter(input_file)
            cv.convert(output_file, start=0, end=None)
            cv.close()
        except Exception as e:
            self.log_message(f"Ошибка при конвертации PDF в DOCX: {str(e)}")
            raise
    
    def pdf_to_image(self, input_file, output_file, fmt):
        self.log_message(f"Конвертация PDF в {fmt.upper()}...")
        try:
            import fitz
            pdf_file = fitz.open(input_file)
            base_name = os.path.splitext(output_file)[0]
            for page_num in range(len(pdf_file)):
                page = pdf_file.load_page(page_num)
                pix = page.get_pixmap()
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                img.save(f"{base_name}_{page_num+1}.{fmt}", fmt.upper())
                
        except Exception as e:
            self.log_message(f"Ошибка при конвертации PDF в изображения: {str(e)}")
            raise
        
    def txt_to_pdf(self, input_file, output_file):
        self.log_message("Конвертация TXT в PDF...")
        try:
            from fpdf import FPDF
            pdf = FPDF()
            pdf.add_page()
            font_path = os.path.join(os.path.dirname(__file__), "DejaVuSans.ttf")
            
            if os.path.exists(font_path):
                pdf.add_font("DejaVu", "", font_path, uni=True)
                pdf.set_font("DejaVu", size=12)
            else:
                raise Exception("Шрифт DejaVuSans.ttf не найден в папке с программой")
            
            with open(input_file, "r", encoding="utf-8") as f:
                for line in f:
                    pdf.cell(0, 10, line.strip(), 0, 1)
            
            pdf.output(output_file)
        except Exception as e:
            raise Exception(f"Ошибка создания PDF: {str(e)}")
    
    def txt_to_docx(self, input_file, output_file):
        self.log_message("Конвертация TXT в DOCX...")
        try:
            doc = Document()
            
            with open(input_file, "r", encoding="utf-8") as f:
                for line in f:
                    doc.add_paragraph(line.strip())
            
            doc.save(output_file)
        except Exception as e:
            raise Exception(f"Ошибка создания DOCX: {str(e)}")
    
    def docx_to_pdf(self, input_file, output_file):
        self.log_message("Конвертация DOCX в PDF...")
        try:
            from docx import Document
            from fpdf import FPDF
            doc = Document(input_file)
            pdf = FPDF()
            pdf.add_page()
            try:
                font_path = os.path.join(os.path.dirname(__file__), "DejaVuSans.ttf")
                if os.path.exists(font_path):
                    pdf.add_font("DejaVu", "", font_path, uni=True)
                    pdf.set_font("DejaVu", size=12)
                else:
                    try:
                        pdf.add_font("ArialUnicode", "", "arialuni.ttf", uni=True)
                        pdf.set_font("ArialUnicode", size=12)
                    except:
                        pdf.set_font("Arial", size=12)
                        self.log_message("Внимание: Используется базовый шрифт, кириллица может не отображаться")
            except Exception as font_error:
                self.log_message(f"Ошибка настройки шрифта: {str(font_error)}")
                raise Exception("Не удалось настроить шрифт с поддержкой кириллицы")

            for element in doc.element.body:
                if element.tag.endswith('p'):
                    para = element
                    text = ' '.join(para.text.split())
                    if text:
                        text_width = pdf.get_string_width(text)
                        page_width = pdf.w - 2*pdf.l_margin
                        
                        if text_width > page_width:
                            pdf.multi_cell(0, 10, text)
                        else:
                            pdf.cell(0, 10, text, ln=1)
                
                elif element.tag.endswith('tbl'):
                    table = element
                    self._process_table(pdf, table)
            
            pdf.output(output_file)
        
        except Exception as e:
            raise Exception(f"Ошибка конвертации DOCX в PDF: {str(e)}")

    def _process_table(self, pdf, table):
        col_count = len(table.xpath('.//w:tr')[0].xpath('.//w:tc'))
        col_width = (pdf.w - 2*pdf.l_margin) / col_count
        
        for row in table.xpath('.//w:tr'):
            for cell in row.xpath('.//w:tc'):
                texts = [t.text for t in cell.xpath('.//w:t')]
                cell_text = ' '.join(''.join(texts).split())
                
                if cell_text:
                    pdf.set_xy(pdf.get_x(), pdf.get_y())
                    pdf.multi_cell(col_width, 10, cell_text, border=1)
            
            pdf.ln()
    
    def docx_to_txt(self, input_file, output_file):
        self.log_message("Конвертация DOCX в TXT...")
        try:
            doc = Document(input_file)
            with open(output_file, "w", encoding="utf-8") as f:
                for para in doc.paragraphs:
                    f.write(para.text + "\n")
        except Exception as e:
            self.log_message(f"Ошибка при конвертации DOCX в TXT: {str(e)}")
            raise
    
    def pptx_to_pdf(self, input_file, output_file):
        self.log_message("Конвертация PPTX в PDF...")
        try:
            prs = Presentation(input_file)
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        pdf.multi_cell(0, 10, shape.text)
                pdf.add_page()
            
            pdf.output(output_file)
            self.log_message("Внимание: Конвертирован только текст из презентации без изображений и форматирования")
        except Exception as e:
            self.log_message(f"Ошибка при конвертации PPTX в PDF: {str(e)}")
            raise
    
    def convert_image(self, input_file, output_file, fmt):
        self.log_message(f"Конвертация изображения в {fmt.upper()}...")
        try:
            img = Image.open(input_file)
            
            if fmt == "jpg":
                img = img.convert("RGB")
            
            save_kwargs = {'quality': 100}
            if fmt in ['png', 'ico']:
                save_kwargs['compress_level'] = 0 
            
            img.save(output_file, fmt.upper(), **save_kwargs)
        except Exception as e:
            self.log_message(f"Ошибка при конвертации изображения: {str(e)}")
            raise
    
    def image_to_pdf(self, input_file, output_file):
        self.log_message("Конвертация изображения в PDF...")
        try:
            with open(output_file, "wb") as f:
                f.write(img2pdf.convert(input_file))
        except Exception as e:
            self.log_message(f"Ошибка при конвертации изображения в PDF: {str(e)}")
            raise
    
    def image_to_ico(self, input_file, output_file, sizes=[(1024,1024), (512,512), (256,256), (128,128), (64,64), (48,48), (32,32), (16,16)]):
        self.log_message("Конвертация в ICO с улучшенным качеством...")
        try:
            img = Image.open(input_file)
            if img.mode != 'RGBA':
                img = img.convert('RGBA')
            
            icon_images = []
            
            for width, height in sizes:
                try:
                    icon = Image.new('RGBA', (width, height), (0, 0, 0, 0))
                    
                    ratio = min(width/img.width, height/img.height)
                    new_width = int(img.width * ratio)
                    new_height = int(img.height * ratio)
                    
                    resized_img = img.resize((new_width, new_height), Image.LANCZOS)
                    
                    x_offset = (width - new_width) // 2
                    y_offset = (height - new_height) // 2
                    
                    icon.paste(resized_img, (x_offset, y_offset), resized_img)
                    icon_images.append(icon)
                    
                except Exception as e:
                    self.log_message(f"Ошибка при создании размера {width}x{height}: {str(e)}")
                    continue
            
            if not icon_images:
                raise Exception("Не удалось создать ни одного варианта иконки")
            
            icon_images[0].save(
                output_file,
                format="ICO",
                sizes=[(i.width, i.height) for i in icon_images],
                quality=100,
                bitmap_format="bmp"
            )
            
            self.log_message(f"Создана ICO с размерами: {', '.join(f'{w}x{h}' for w, h in sizes)}")
            
        except Exception as e:
            self.log_message(f"Ошибка при конвертации в ICO: {str(e)}")
            raise
    
    def ico_to_image(self, input_file, output_file, fmt):
        self.log_message(f"Конвертация ICO в {fmt.upper()}...")
        try:
            img = Image.open(input_file)
            
            if fmt == "jpg":
                img = img.convert("RGB")
            
            img.save(output_file, fmt.upper())
        except Exception as e:
            self.log_message(f"Ошибка при конвертации ICO: {str(e)}")
            raise


if __name__ == "__main__":
    root = tk.Tk()
    app = FileConverterApp(root)
    root.mainloop()
